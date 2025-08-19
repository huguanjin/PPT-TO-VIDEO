"""
步骤1: PPT解析器
解析PPT文件，提取每页内容和备注信息
"""
import os
from pathlib import Path
from typing import Dict, Any, List, Optional, Callable
from datetime import datetime
import asyncio
import logging

from pptx import Presentation
from pptx.slide import Slide
from PIL import Image, ImageDraw, ImageFont
import io

from utils.logger import get_logger
from utils.file_manager import FileManager

# 尝试导入Windows COM接口
try:
    import win32com.client
    import pythoncom
    WIN32_AVAILABLE = True
except ImportError:
    WIN32_AVAILABLE = False


class PPTParser:
    """PPT解析器"""
    
    def __init__(self, project_dir: Path):
        self.project_dir = Path(project_dir)
        self.file_manager = FileManager(project_dir)
        self.logger = get_logger(__name__, self.project_dir / "logs")
        self._current_ppt_file = None
        
    async def parse_ppt(self, ppt_file_path: str, progress_callback: Optional[Callable[[int], None]] = None) -> Dict[str, Any]:
        """
        解析PPT文件或加载现有项目数据
        
        Args:
            ppt_file_path: PPT文件路径（对于现有项目可以为None）
            progress_callback: 进度回调函数
            
        Returns:
            解析结果字典
        """
        try:
            # 首先检查是否存在现有的幻灯片数据
            slides_dir = self.project_dir / "slides"
            metadata_file = slides_dir / "slides_metadata.json"
            
            if slides_dir.exists() and metadata_file.exists():
                self.logger.info("发现现有幻灯片数据，尝试加载...")
                existing_data = await self._load_existing_slides_data()
                
                if existing_data and existing_data.get("slides"):
                    self.logger.info(f"成功加载现有数据: {len(existing_data['slides'])} 张幻灯片")
                    if progress_callback:
                        progress_callback(100)
                    return existing_data
                else:
                    self.logger.warning("现有数据格式不正确，将重新解析")
            
            # 如果没有现有数据或数据无效，则进行PPT解析
            if not ppt_file_path:
                raise ValueError("没有找到现有幻灯片数据，且未提供PPT文件路径")
                
            self.logger.info(f"开始解析PPT文件: {ppt_file_path}")
            
            # 保存当前PPT文件路径以供COM接口使用
            self._current_ppt_file = ppt_file_path
            
            # 使用python-pptx解析PPT
            presentation = Presentation(ppt_file_path)
            slides_info = []
            
            total_slides = len(presentation.slides)
            self.logger.info(f"PPT总共有 {total_slides} 张幻灯片")
            
            for i, slide in enumerate(presentation.slides):
                slide_number = i + 1
                self.logger.info(f"处理第 {slide_number} 张幻灯片")
                
                # 解析单张幻灯片
                slide_info = await self._parse_slide(slide, slide_number)
                slides_info.append(slide_info)
                
                # 更新进度
                if progress_callback:
                    progress = int((i + 1) / total_slides * 100)
                    progress_callback(progress)
                
                # 让出控制权
                await asyncio.sleep(0.01)
            
            # 计算总时长
            total_duration = sum(
                self._estimate_duration(slide_info.get("notes", ""))
                for slide_info in slides_info
            )
            
            result = {
                "ppt_file": ppt_file_path,
                "total_slides": total_slides,
                "slides": slides_info,
                "total_duration": total_duration,
                "parsed_at": datetime.now().isoformat()
            }
            
            self.logger.info(f"PPT解析完成，共 {total_slides} 张幻灯片，预计总时长 {total_duration:.1f} 秒")
            return result
            
        except Exception as e:
            self.logger.error(f"解析PPT失败: {e}")
            raise
    
    async def _load_existing_slides_data(self) -> Optional[Dict[str, Any]]:
        """
        加载现有的幻灯片数据
        
        Returns:
            现有的幻灯片数据字典，如果加载失败则返回None
        """
        try:
            slides_dir = self.project_dir / "slides"
            metadata_file = slides_dir / "slides_metadata.json"
            
            if not metadata_file.exists():
                self.logger.warning("幻灯片元数据文件不存在")
                return None
            
            # 读取元数据文件
            import json
            with open(metadata_file, 'r', encoding='utf-8') as f:
                metadata = json.load(f)
            
            slides_info = []
            slide_files = list(slides_dir.glob("slide_*.png"))
            
            if not slide_files:
                self.logger.warning("没有找到幻灯片图片文件")
                return None
            
            self.logger.info(f"找到 {len(slide_files)} 个幻灯片图片文件")
            
            # 根据现有图片文件重建幻灯片信息
            slide_files.sort()  # 按文件名排序
            
            for i, slide_file in enumerate(slide_files, 1):
                slide_number = i
                
                # 从元数据中获取对应的幻灯片信息
                slide_metadata = None
                if isinstance(metadata, dict) and "slides" in metadata:
                    slides_list = metadata["slides"]
                    for slide in slides_list:
                        if slide.get("slide_number") == slide_number:
                            slide_metadata = slide
                            break
                
                # 构建幻灯片信息
                slide_info = {
                    "slide_id": f"{slide_number:03d}",
                    "slide_number": slide_number,
                    "title": slide_metadata.get("title", f"幻灯片{slide_number}") if slide_metadata else f"幻灯片{slide_number}",
                    "image_file": f"slides/{slide_file.name}",
                    "notes": slide_metadata.get("remark", "") if slide_metadata else "",
                    "notes_word_count": len(slide_metadata.get("remark", "")) if slide_metadata else 0,
                    "extracted_at": datetime.now().isoformat(),
                    "loaded_from_existing": True  # 标记为从现有数据加载
                }
                
                slides_info.append(slide_info)
                self.logger.info(f"加载第 {slide_number} 张幻灯片: {slide_info['title']}")
            
            # 计算总时长
            total_duration = sum(
                self._estimate_duration(slide_info.get("notes", ""))
                for slide_info in slides_info
            )
            
            result = {
                "ppt_file": None,  # 现有项目没有原始PPT文件
                "total_slides": len(slides_info),
                "slides": slides_info,
                "total_duration": total_duration,
                "parsed_at": datetime.now().isoformat(),
                "loaded_from_existing": True,
                "original_metadata": metadata  # 保存原始元数据
            }
            
            self.logger.info(f"成功加载现有幻灯片数据，共 {len(slides_info)} 张幻灯片，预计总时长 {total_duration:.1f} 秒")
            return result
            
        except Exception as e:
            self.logger.error(f"加载现有幻灯片数据失败: {e}")
            return None
    
    async def _parse_slide(self, slide: Slide, slide_number: int) -> Dict[str, Any]:
        """解析单张幻灯片"""
        # 提取标题
        title = self._extract_title(slide)
        
        # 保存幻灯片图片
        image_file = await self._save_slide_image(slide, slide_number)
        
        # 提取备注
        notes = self._extract_notes(slide)
        
        slide_info = {
            "slide_id": f"{slide_number:03d}",
            "slide_number": slide_number,
            "title": title,
            "image_file": image_file,
            "notes": notes,
            "notes_word_count": len(notes) if notes else 0,
            "extracted_at": datetime.now().isoformat()
        }
        
        return slide_info
    
    def _extract_title(self, slide: Slide) -> str:
        """提取幻灯片标题"""
        try:
            if hasattr(slide, 'shapes'):
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # 通常第一个有文字的形状是标题
                        return shape.text.strip()
            return f"幻灯片标题"
        except Exception as e:
            self.logger.warning(f"提取标题失败: {e}")
            return "无标题"
    
    def _extract_notes(self, slide: Slide) -> Optional[str]:
        """提取幻灯片备注"""
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    return notes_text if notes_text else None
            return None
        except Exception as e:
            self.logger.warning(f"提取备注失败: {e}")
            return None
    
    async def _save_slide_image(self, slide: Slide, slide_number: int) -> str:
        """
        保存幻灯片为图片
        支持多种方法：Windows COM、占位符图片
        """
        image_filename = f"slide_{slide_number:03d}.png"
        image_path = self.file_manager.slides_dir / image_filename
        
        try:
            # 确保目录存在
            image_path.parent.mkdir(parents=True, exist_ok=True)
            
            # 启用COM接口测试真实PPT图片导出
            com_enabled = True
            
            # 优先尝试使用Windows COM接口导出真实幻灯片图片
            if WIN32_AVAILABLE and self._current_ppt_file and com_enabled:
                try:
                    success = self._export_slide_with_com_sync(slide_number, image_path)
                    if success and image_path.exists():
                        self.logger.info(f"使用COM接口成功导出幻灯片 {slide_number}")
                        return image_filename
                    else:
                        self.logger.warning(f"COM接口导出失败，回退到占位符图片")
                except Exception as e:
                    self.logger.warning(f"COM接口导出异常: {e}, 回退到占位符图片")
            
            # 降级到创建高质量占位符图片
            self._create_placeholder_image(slide, image_path)
            self.logger.info(f"创建占位符图片: {image_filename}")
            return image_filename
            
        except Exception as e:
            self.logger.error(f"保存幻灯片图片失败: {e}")
            # 创建简单的占位符
            self._create_simple_placeholder(slide_number, image_path)
            return image_filename
    
    def _export_slide_with_com_sync(self, slide_number: int, output_path: Path) -> bool:
        """
        同步方式使用Windows COM接口导出真实的幻灯片图片
        """
        if not WIN32_AVAILABLE:
            return False
        
        ppt_app = None
        presentation = None
        try:
            # 初始化COM
            pythoncom.CoInitialize()
            
            # 创建PowerPoint应用对象
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = True  # 必须设置为True，否则会出错
            
            # 打开PPT文件，确保路径编码正确
            abs_ppt_path = str(Path(self._current_ppt_file).absolute())
            
            # 处理中文路径，确保COM接口能够正确处理
            try:
                # 检查路径是否包含中文字符
                abs_ppt_path.encode('ascii')
            except UnicodeEncodeError:
                # 如果包含非ASCII字符，尝试使用短路径名
                try:
                    import win32api
                    abs_ppt_path = win32api.GetShortPathName(abs_ppt_path)
                    self.logger.info(f"使用短路径名处理中文路径: {abs_ppt_path}")
                except:
                    self.logger.warning("无法获取短路径名，使用原始路径")
            
            presentation = ppt_app.Presentations.Open(abs_ppt_path)
            
            # 导出特定幻灯片，确保输出路径编码正确
            slide = presentation.Slides(slide_number)  # PowerPoint使用1索引
            abs_output_path = str(output_path.absolute())
            
            # 确保输出路径不包含中文字符，避免COM接口问题
            try:
                abs_output_path.encode('ascii')
            except UnicodeEncodeError:
                # 如果输出路径包含中文，使用短路径名
                try:
                    import win32api
                    parent_dir = str(output_path.parent.absolute())
                    short_parent = win32api.GetShortPathName(parent_dir)
                    abs_output_path = os.path.join(short_parent, output_path.name)
                    self.logger.info(f"使用短路径名处理输出路径: {abs_output_path}")
                except:
                    self.logger.warning("无法获取输出目录短路径名，使用原始路径")
            
            slide.Export(abs_output_path, "PNG", 1920, 1080)  # 导出为高分辨率PNG
            
            return True
            
        except Exception as e:
            self.logger.error(f"COM导出操作失败: {e}")
            return False
        finally:
            # 清理资源
            try:
                if presentation:
                    presentation.Close()
                if ppt_app:
                    ppt_app.Quit()
                pythoncom.CoUninitialize()
            except:
                pass
    
    def _create_placeholder_image(self, slide: Slide, output_path: Path):
        """
        创建高质量占位符图片，包含幻灯片的文本内容和背景
        """
        try:
            # 创建1920x1080的高分辨率图片
            width, height = 1920, 1080
            
            # 尝试提取幻灯片背景
            background_color = self._extract_slide_background(slide)
            image = Image.new('RGB', (width, height), color=background_color)
            draw = ImageDraw.Draw(image)
            
            # 如果是白色背景，添加一些设计元素
            if background_color == (255, 255, 255):  # 白色背景
                # 添加渐变效果
                for i in range(height):
                    gradient_value = int(240 + (i / height) * 15)  # 从240到255的渐变
                    color = (gradient_value, gradient_value, gradient_value)
                    draw.line([(0, i), (width, i)], fill=color)
                
                # 添加装饰性边框
                border_color = (200, 200, 200)
                border_width = 8
                draw.rectangle([border_width//2, border_width//2, width-border_width//2, height-border_width//2], 
                             outline=border_color, width=border_width)
                
                # 添加顶部装饰条
                header_height = 80
                header_color = (50, 120, 200)  # 蓝色装饰条
                draw.rectangle([0, 0, width, header_height], fill=header_color)
            else:
                # 非白色背景，添加简单边框
                border_color = tuple(max(0, c - 50) for c in background_color)  # 稍微深一点的颜色做边框
                border_width = 4
                draw.rectangle([border_width//2, border_width//2, width-border_width//2, height-border_width//2], 
                             outline=border_color, width=border_width)
            
            # 提取幻灯片内容
            title = self._extract_slide_title(slide) or "幻灯片标题"
            content = self._extract_slide_text(slide) or "幻灯片内容"
            
            # 尝试加载中文字体
            try:
                # 优先使用中文字体
                title_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 60)  # 微软雅黑
                content_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 36)
            except:
                try:
                    title_font = ImageFont.truetype("arial.ttf", 60)
                    content_font = ImageFont.truetype("arial.ttf", 36)
                except:
                    # 使用默认字体
                    title_font = ImageFont.load_default()
                    content_font = ImageFont.load_default()
            
            # 根据背景颜色选择文字颜色
            text_color = self._get_contrasting_color(background_color)
            title_color = text_color
            
            # 绘制标题
            title_y = 150 if background_color == (255, 255, 255) else 120  # 白色背景时为装饰条留空间
            self._draw_wrapped_text(draw, title, title_font, title_color, 100, title_y, width-200, 150)
            
            # 绘制内容
            content_y = 320
            max_content_height = height - content_y - 100
            self._draw_wrapped_text(draw, content, content_font, text_color, 100, content_y, width-200, max_content_height)
            
            # 在右下角添加水印
            watermark = "PPT转视频工具生成"
            try:
                watermark_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 24)
            except:
                try:
                    watermark_font = ImageFont.truetype("arial.ttf", 24)
                except:
                    watermark_font = ImageFont.load_default()
            
            watermark_color = tuple(min(255, c + 100) for c in text_color) if sum(text_color) < 300 else tuple(max(0, c - 100) for c in text_color)
            watermark_bbox = draw.textbbox((0, 0), watermark, font=watermark_font)
            watermark_width = watermark_bbox[2] - watermark_bbox[0]
            watermark_x = width - watermark_width - 30
            watermark_y = height - 50
            draw.text((watermark_x, watermark_y), watermark, fill=watermark_color, font=watermark_font)
            
            # 保存图片
            image.save(output_path, 'PNG', quality=95)
            
        except Exception as e:
            self.logger.error(f"创建占位符图片失败: {e}")
            # 创建简单占位符
            self._create_simple_placeholder(1, output_path)
    
    def _extract_slide_background(self, slide: Slide) -> tuple:
        """提取幻灯片背景颜色"""
        try:
            # 尝试获取幻灯片背景
            if hasattr(slide, 'background'):
                background = slide.background
                if hasattr(background, 'fill') and hasattr(background.fill, 'fore_color'):
                    try:
                        color = background.fill.fore_color.rgb
                        if color:
                            # 处理RGBColor对象，转换为RGB元组
                            if hasattr(color, '__int__'):
                                rgb_int = int(color)
                                return (rgb_int >> 16 & 0xFF, rgb_int >> 8 & 0xFF, rgb_int & 0xFF)
                            elif hasattr(color, 'rgb'):
                                rgb_int = int(color.rgb)
                                return (rgb_int >> 16 & 0xFF, rgb_int >> 8 & 0xFF, rgb_int & 0xFF)
                    except Exception as e:
                        self.logger.debug(f"处理背景颜色失败: {e}")
            
            # 如果无法获取背景，检查是否有背景形状
            if hasattr(slide, 'shapes'):
                for shape in slide.shapes:
                    if hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color'):
                        try:
                            color = shape.fill.fore_color.rgb
                            if color:
                                if hasattr(color, '__int__'):
                                    rgb_int = int(color)
                                    if rgb_int != 0:  # 不是黑色
                                        return (rgb_int >> 16 & 0xFF, rgb_int >> 8 & 0xFF, rgb_int & 0xFF)
                                elif hasattr(color, 'rgb'):
                                    rgb_int = int(color.rgb)
                                    if rgb_int != 0:  # 不是黑色
                                        return (rgb_int >> 16 & 0xFF, rgb_int >> 8 & 0xFF, rgb_int & 0xFF)
                        except Exception as e:
                            self.logger.debug(f"处理形状颜色失败: {e}")
                            continue
            
            # 默认返回白色
            return (255, 255, 255)
            
        except Exception as e:
            self.logger.warning(f"提取背景颜色失败: {e}")
            return (255, 255, 255)  # 默认白色
    
    def _get_contrasting_color(self, background_color: tuple) -> tuple:
        """根据背景颜色获取对比色文字颜色"""
        # 计算背景颜色的亮度
        brightness = sum(background_color) / 3
        
        if brightness > 128:
            # 浅色背景使用深色文字
            return (50, 50, 50)
        else:
            # 深色背景使用浅色文字
            return (240, 240, 240)
    
    def _draw_wrapped_text(self, draw, text: str, font, color, x: int, y: int, max_width: int, max_height: int):
        """绘制自动换行的文本"""
        try:
            words = text.split()
            lines = []
            current_line = []
            
            for word in words:
                test_line = ' '.join(current_line + [word])
                bbox = draw.textbbox((0, 0), test_line, font=font)
                text_width = bbox[2] - bbox[0]
                
                if text_width <= max_width:
                    current_line.append(word)
                else:
                    if current_line:
                        lines.append(' '.join(current_line))
                        current_line = [word]
                    else:
                        lines.append(word)  # 单词太长，强制换行
            
            if current_line:
                lines.append(' '.join(current_line))
            
            # 绘制文本行
            line_height = 50 if font == draw._font else 40
            current_y = y
            
            for line in lines:
                if current_y + line_height > y + max_height:
                    break  # 超出最大高度
                draw.text((x, current_y), line, fill=color, font=font)
                current_y += line_height
                
        except Exception as e:
            # 如果出错，直接绘制原始文本
            draw.text((x, y), text[:100] + "..." if len(text) > 100 else text, fill=color, font=font)
    
    def _create_simple_placeholder(self, slide_number: int, output_path: Path):
        """创建简单的占位符图片"""
        try:
            width, height = 1920, 1080
            image = Image.new('RGB', (width, height), color='lightgray')
            draw = ImageDraw.Draw(image)
            
            # 绘制文本
            text = f"幻灯片 {slide_number}"
            
            # 计算文本位置（居中）
            try:
                font = ImageFont.truetype("arial.ttf", 72)
            except:
                font = ImageFont.load_default()
            
            bbox = draw.textbbox((0, 0), text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            
            x = (width - text_width) // 2
            y = (height - text_height) // 2
            
            draw.text((x, y), text, fill='black', font=font)
            
            # 保存图片
            image.save(output_path, 'PNG')
            
        except Exception as e:
            self.logger.error(f"创建简单占位符失败: {e}")
            # 最后的备选方案
            simple_image = Image.new('RGB', (1920, 1080), color='white')
            simple_image.save(output_path, 'PNG')
    
    def _estimate_duration(self, text: str) -> float:
        """
        估算文本的朗读时长（秒）
        """
        if not text or not text.strip():
            return 3.0  # 默认3秒
        
        # 中文平均语速约为 3-4 字/秒，这里使用3.5字/秒
        chars_per_second = 3.5
        char_count = len(text.replace(" ", "").replace("\n", ""))
        
        return char_count / chars_per_second
    
    def _extract_slide_title(self, slide: Slide) -> Optional[str]:
        """
        提取幻灯片标题（用于占位符图片）
        """
        try:
            if hasattr(slide, 'shapes'):
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # 返回第一个非空文本作为标题
                        text = shape.text.strip()
                        if len(text) <= 50:  # 标题通常较短
                            return text
                        else:
                            # 如果文本太长，取前50个字符作为标题
                            return text[:50] + "..."
            return None
        except Exception as e:
            self.logger.warning(f"提取幻灯片标题失败: {e}")
            return None
    
    def _extract_slide_text(self, slide: Slide) -> Optional[str]:
        """
        提取幻灯片所有文本内容（用于占位符图片）
        """
        try:
            all_text = []
            if hasattr(slide, 'shapes'):
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        all_text.append(shape.text.strip())
            
            return "\n".join(all_text) if all_text else None
        except Exception as e:
            self.logger.warning(f"提取幻灯片文本失败: {e}")
            return None
