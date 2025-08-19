"""
步骤1: PPT解析器
解析PPT文件，提取每页内容和备注信息
"""
import os
from pathlib import Path
from typing import Dict, Any, List, Optional, Callable
from datetime import datetime
import asyncio
import logging

from pptx import Presentation
from pptx.slide import Slide
from PIL import Image
import io

from utils.logger import get_logger
from utils.file_manager import FileManager

# 尝试导入Windows COM接口
try:
    import win32com.client
    import pythoncom
    WIN32_AVAILABLE = True
except ImportError:
    WIN32_AVAILABLE = False
                    if hasattr(shape, 'text') and shape.text.strip():
                        # 返回第一个非空文本作为标题
                        text = shape.text.strip()
                        if len(text) <= 50:  # 标题通常较短
                            return text
                        else:
                            # 如果文本太长，取前50个字符作为标题
                            return text[:50] + "..."
            return None
        except Exception as e:
            self.logger.warning(f"提取幻灯片标题失败: {e}")
            return None
    
    def _extract_slide_text(self, slide: Slide) -> Optional[str]:
        """
        提取幻灯片所有文本内容（用于占位符图片）
        """
        try:
            all_text = []
            if hasattr(slide, 'shapes'):
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        all_text.append(shape.text.strip())
            
            return "\n".join(all_text) if all_text else None
        except Exception as e:
            self.logger.warning(f"提取幻灯片文本失败: {e}")
            return Noneport FileManager

class PPTParser:
    """PPT解析器"""
    
    def __init__(self, project_dir: Path):
        self.project_dir = Path(project_dir)
        self.file_manager = FileManager(project_dir)
        self.logger = get_logger(__name__, self.project_dir / "logs")
        
    async def parse_ppt(self, ppt_file_path: str, progress_callback: Optional[Callable[[int], None]] = None) -> Dict[str, Any]:
        """
        解析PPT文件
        
        Args:
            ppt_file_path: PPT文件路径
            progress_callback: 进度回调函数
            
        Returns:
            解析结果字典
        """
        try:
            self.logger.info(f"开始解析PPT文件: {ppt_file_path}")
            
            # 保存当前PPT文件路径以供COM接口使用
            self._current_ppt_file = ppt_file_path
            
            # 确保必要的目录存在
            self.file_manager.create_directory_structure()
            
            # 打开PPT文件
            presentation = Presentation(ppt_file_path)
            total_slides = len(presentation.slides)
            
            self.logger.info(f"PPT总页数: {total_slides}")
            
            slides_data = {
                "total_slides": total_slides,
                "parsing_completed": False,
                "parsing_timestamp": datetime.now().isoformat(),
                "source_file": ppt_file_path,
                "slides": []
            }
            
            # 解析每一页幻灯片
            for i, slide in enumerate(presentation.slides):
                if progress_callback:
                    progress = int((i / total_slides) * 100)
                    progress_callback(progress)
                
                slide_info = await self._parse_slide(slide, i + 1)
                slides_data["slides"].append(slide_info)
                
                self.logger.info(f"解析完成第 {i + 1} 页")
                
                # 模拟异步处理
                await asyncio.sleep(0.1)
            
            slides_data["parsing_completed"] = True
            
            # 保存解析结果
            self.file_manager.save_slides_metadata(slides_data)
            
            # 保存讲话稿数据到统一的JSON文件
            scripts_data = self._generate_scripts_metadata(slides_data)
            self.file_manager.save_scripts_metadata(scripts_data)
            
            if progress_callback:
                progress_callback(100)
                
            self.logger.info("PPT解析完成")
            return slides_data
            
        except Exception as e:
            self.logger.error(f"PPT解析失败: {e}", exc_info=True)
            raise
    
    async def _parse_slide(self, slide: Slide, slide_number: int) -> Dict[str, Any]:
        """
        解析单页幻灯片
        
        Args:
            slide: 幻灯片对象
            slide_number: 幻灯片编号
            
        Returns:
            幻灯片信息字典
        """
        # 提取幻灯片标题
        title = self._extract_title(slide)
        
        # 提取备注信息
        notes = self._extract_notes(slide)
        
        # 生成幻灯片图片
        image_file = await self._save_slide_image(slide, slide_number)
        
        slide_info = {
            "slide_id": f"{slide_number:03d}",
            "slide_number": slide_number,
            "title": title,
            "image_file": image_file,
            "notes": notes,
            "notes_word_count": len(notes) if notes else 0,
            "extracted_at": datetime.now().isoformat()
        }
        
        return slide_info
    
    def _extract_title(self, slide: Slide) -> str:
        """提取幻灯片标题"""
        try:
            if hasattr(slide, 'shapes'):
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # 通常第一个有文字的形状是标题
                        return shape.text.strip()
            return f"幻灯片标题"
        except Exception as e:
            self.logger.warning(f"提取标题失败: {e}")
            return "无标题"
    
    def _extract_notes(self, slide: Slide) -> Optional[str]:
        """提取幻灯片备注"""
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    return notes_text if notes_text else None
            return None
        except Exception as e:
            self.logger.warning(f"提取备注失败: {e}")
            return None
    
    async def _save_slide_image(self, slide: Slide, slide_number: int) -> str:
        """
        保存幻灯片为图片
        支持多种方法：Windows COM、占位符图片
        """
        image_filename = f"slide_{slide_number:03d}.png"
        image_path = self.file_manager.slides_dir / image_filename
        
        try:
            # 确保目录存在
            image_path.parent.mkdir(parents=True, exist_ok=True)
            
            # 尝试使用Windows COM方式导出真实的PPT图片
            if await self._export_slide_with_com(slide_number, image_path):
                self.logger.info(f"保存幻灯片图片(COM): {image_path}")
                return image_filename
            
            # 如果COM方式失败，创建一个更详细的占位符图片
            self.logger.warning("COM方式失败，使用占位符图片")
            await self._create_placeholder_image(slide, slide_number, image_path)
            self.logger.info(f"保存幻灯片图片(占位符): {image_path}")
            
        except Exception as e:
            self.logger.error(f"保存幻灯片图片失败: {e}")
            # 即使图片生成失败，也不影响其他功能
            # 创建一个简单的空白图片
            try:
                from PIL import Image
                img = Image.new('RGB', (1920, 1080), color='lightgray')
                img.save(image_path)
                self.logger.warning(f"使用简单占位符图片: {image_path}")
            except:
                pass
        
        return image_filename
    
    async def _export_slide_with_com(self, slide_number: int, output_path: Path) -> bool:
        """
        使用Windows COM接口导出PPT幻灯片为图片
        """
        try:
            import win32com.client
            import os
            import time
            
            # 创建PowerPoint应用实例
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = False
            
            # 打开当前正在处理的PPT文件
            # 注意：这里需要传递当前PPT文件的完整路径
            ppt_file = getattr(self, '_current_ppt_file', None)
            if not ppt_file:
                return False
            
            presentation = ppt_app.Presentations.Open(os.path.abspath(ppt_file))
            
            # 获取指定幻灯片
            slide = presentation.Slides(slide_number)  # COM接口中slide索引从1开始
            
            # 导出为图片
            slide.Export(str(output_path), "PNG", 1920, 1080)
            
            # 清理资源
            presentation.Close()
            ppt_app.Quit()
            
            # 检查文件是否成功创建
            if output_path.exists():
                return True
            
        except Exception as e:
            self.logger.warning(f"COM导出失败: {e}")
            try:
                # 清理可能残留的COM对象
                if 'ppt_app' in locals():
                    ppt_app.Quit()
            except:
                pass
        
        return False
    
    async def _create_placeholder_image(self, slide: Slide, slide_number: int, image_path: Path):
        """
        创建更详细的占位符图片，包含幻灯片基本信息
        """
        from PIL import Image, ImageDraw, ImageFont
        
        img = Image.new('RGB', (1920, 1080), color='white')
        draw = ImageDraw.Draw(img)
        
        # 添加边框
        draw.rectangle([50, 50, 1870, 1030], outline='gray', width=3)
        
        # 尝试获取幻灯片标题
        title = self._extract_slide_title(slide) or f"第 {slide_number} 页"
        
        # 尝试获取幻灯片文本内容
        content = self._extract_slide_text(slide)
        
        try:
            # 尝试使用中文字体
            title_font = ImageFont.truetype("msyh.ttc", 60)  # 微软雅黑
            content_font = ImageFont.truetype("msyh.ttc", 32)
        except:
            try:
                # 备选字体
                title_font = ImageFont.truetype("arial.ttf", 60)
                content_font = ImageFont.truetype("arial.ttf", 32)
            except:
                # 默认字体
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
        
        # 绘制标题
        title_bbox = draw.textbbox((0, 0), title, font=title_font)
        title_width = title_bbox[2] - title_bbox[0]
        title_x = (1920 - title_width) // 2
        draw.text((title_x, 150), title, fill='black', font=title_font)
        
        # 绘制页码
        page_text = f"页码: {slide_number}"
        draw.text((100, 100), page_text, fill='gray', font=content_font)
        
        # 绘制内容预览
        if content:
            # 限制内容长度
            preview_content = content[:200] + "..." if len(content) > 200 else content
            # 简单换行处理
            lines = []
            words = preview_content.split()
            current_line = ""
            for word in words:
                test_line = current_line + " " + word if current_line else word
                bbox = draw.textbbox((0, 0), test_line, font=content_font)
                line_width = bbox[2] - bbox[0]
                if line_width < 1600:  # 留出边距
                    current_line = test_line
                else:
                    if current_line:
                        lines.append(current_line)
                        current_line = word
                    if len(lines) >= 20:  # 最多20行
                        break
            if current_line:
                lines.append(current_line)
            
            # 绘制文本行
            y_offset = 250
            for line in lines:
                draw.text((150, y_offset), line, fill='darkblue', font=content_font)
                y_offset += 40
        
        # 添加水印
        watermark = "PPT预览图 - 实际内容请参考原PPT文件"
        draw.text((150, 950), watermark, fill='lightgray', font=content_font)
        
        # 保存图片
        img.save(image_path)
    
    def _generate_scripts_metadata(self, slides_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成讲话稿元数据
        
        Args:
            slides_data: 幻灯片数据
            
        Returns:
            讲话稿元数据字典
        """
        scripts_data = {
            "total_scripts": slides_data["total_slides"],
            "extraction_completed": True,
            "extraction_timestamp": datetime.now().isoformat(),
            "total_word_count": sum(slide.get("notes_word_count", 0) for slide in slides_data["slides"]),
            "scripts": []
        }
        
        for slide in slides_data["slides"]:
            script_info = {
                "script_id": slide["slide_id"],
                "slide_number": slide["slide_number"],
                "slide_title": slide["title"],
                "script_content": slide["notes"] or "",
                "word_count": slide["notes_word_count"],
                "estimated_duration_seconds": self._estimate_speech_duration(slide["notes"]),
                "extracted_at": slide["extracted_at"]
            }
            scripts_data["scripts"].append(script_info)
        
        return scripts_data
    
    def _estimate_speech_duration(self, text: str) -> float:
        """
        估算语音时长（基于字数和平均语速）
        
        Args:
            text: 文本内容
            
        Returns:
            估算的语音时长（秒）
        """
        if not text:
            return 0.0
        
        # 中文平均语速约为 3-4 字/秒，这里使用3.5字/秒
        chars_per_second = 3.5
        char_count = len(text.replace(" ", "").replace("\n", ""))
        
        return round(char_count / chars_per_second, 2)
    
    async def get_slide_count(self, ppt_file_path: str) -> int:
        """获取PPT页数"""
        try:
            presentation = Presentation(ppt_file_path)
            return len(presentation.slides)
        except Exception as e:
            self.logger.error(f"获取PPT页数失败: {e}")
            return 0
